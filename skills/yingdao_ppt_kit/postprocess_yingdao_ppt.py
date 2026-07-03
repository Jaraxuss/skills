"""
Legacy Yingdao-style PPT Postprocessor
--------------------------------------
Apply reusable visual rules to an existing .pptx file:
- all visible text uses Microsoft YaHei
- code blocks use 12pt and light background
- table header fill becomes red/pink, header text black
- top-right logo can be enlarged while keeping top-right anchor

This script is a compatibility helper only. It cannot replace AI-led slide
planning, visual asset selection, rendered QA, and manual-quality revision for
final customer-facing decks.

Usage:
    python postprocess_yingdao_ppt.py input.pptx output.pptx --logo-scale 2
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

FONT = "Microsoft YaHei"

COLORS = {
    "black": RGBColor(0x1F, 0x1F, 0x1F),
    "body": RGBColor(0x33, 0x33, 0x33),
    "muted": RGBColor(0x77, 0x77, 0x77),
    "red": RGBColor(0xF2, 0x3A, 0x4A),
    "dark_red": RGBColor(0xB8, 0x3A, 0x4A),
    "table_header": RGBColor(0xF1, 0x9B, 0xAA),
    "pink": RGBColor(0xFF, 0xF3, 0xF5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "code_bg": RGBColor(0xFF, 0xF6, 0xF7),
    "code_border": RGBColor(0xEE, 0xA3, 0xAE),
    "code_text": RGBColor(0x40, 0x36, 0x3A),
    "code_keyword": RGBColor(0xCF, 0x2F, 0x43),
    "code_func": RGBColor(0x7E, 0x3F, 0x98),
    "code_string": RGBColor(0xB9, 0x59, 0x00),
    "code_comment": RGBColor(0x7B, 0x7B, 0x7B),
}

KEYWORDS = {
    "import", "from", "as", "for", "if", "elif", "else", "return", "in", "not",
    "and", "or", "True", "False", "None", "with", "def", "class", "lambda",
}

CODE_HINTS = (
    "import pandas", "pd.", "df[", ".to_excel", ".read_excel", ".merge", ".groupby",
    "drop_duplicates", "dropna", "fillna", "to_datetime", "to_numeric",
)


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            yield shape


def set_all_fonts(prs: Presentation, min_body_size: int = 14) -> None:
    """Set all text to Microsoft YaHei and lift tiny non-footer body text."""
    for slide in prs.slides:
        for shape in iter_text_shapes(slide):
            text = shape.text.strip()
            is_footer = text == "From human doing to human being." or (text.isdigit() and len(text) <= 2)
            is_code = looks_like_code(text)
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue
                    run.font.name = FONT
                    if is_code:
                        run.font.size = Pt(12)
                    elif is_footer:
                        run.font.size = Pt(8)
                    elif run.font.size is None or run.font.size.pt < min_body_size:
                        run.font.size = Pt(min_body_size)
                    if run.font.color.type is None:
                        run.font.color.rgb = COLORS["body"]


def looks_like_code(text: str) -> bool:
    return any(hint in text for hint in CODE_HINTS) or text.strip().startswith("# 建议输出")


def restyle_tables(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            table = shape.table
            for r, row in enumerate(table.rows):
                for c, cell in enumerate(row.cells):
                    cell.fill.solid()
                    if r == 0:
                        cell.fill.fore_color.rgb = COLORS["table_header"]
                    else:
                        cell.fill.fore_color.rgb = COLORS["white"] if r % 2 else COLORS["pink"]
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = FONT
                            run.font.size = Pt(14)
                            run.font.color.rgb = COLORS["black"] if r == 0 else COLORS["body"]
                            if r == 0:
                                run.font.bold = True


def restyle_manual_table_headers(prs: Presentation) -> None:
    """Handle table-like layouts made from rectangles and text boxes.

    This detects black/dark header rectangles near the first content row and turns
    them into red/pink header blocks. It is intentionally conservative.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if shape.fill.type != MSO_FILL.SOLID:
                continue
            rgb = shape.fill.fore_color.rgb
            if rgb is None:
                continue
            is_dark = (rgb[0] < 55 and rgb[1] < 55 and rgb[2] < 55)
            is_header_band = Inches(1.0) <= shape.top <= Inches(1.8) and shape.width > Inches(0.5)
            if is_dark and is_header_band:
                shape.fill.solid()
                shape.fill.fore_color.rgb = COLORS["table_header"]
                shape.line.color.rgb = COLORS["table_header"]


def restyle_code_blocks(prs: Presentation) -> None:
    for slide in prs.slides:
        shapes = list(slide.shapes)
        for i, shape in enumerate(shapes):
            # background rectangles: old black or explicitly named code background
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.fill.type == MSO_FILL.SOLID:
                rgb = shape.fill.fore_color.rgb
                is_old_dark_code_bg = rgb is not None and rgb[0] < 60 and rgb[1] < 60 and rgb[2] < 60
                is_named_code_bg = getattr(shape, "name", "") == "CODE_BACKGROUND"
                if is_old_dark_code_bg or is_named_code_bg:
                    # only apply if the next/nearby text shape looks like code
                    nearby_text = ""
                    for j in range(i + 1, min(i + 3, len(shapes))):
                        if hasattr(shapes[j], "text"):
                            nearby_text += shapes[j].text
                    if is_named_code_bg or looks_like_code(nearby_text):
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = COLORS["code_bg"]
                        shape.line.color.rgb = COLORS["code_border"]
                        shape.line.width = Pt(0.8)
        for shape in iter_text_shapes(slide):
            if looks_like_code(shape.text):
                apply_code_highlight(shape)


def apply_code_highlight(shape) -> None:
    original = shape.text
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(original.strip("\n").split("\n")):
        paragraph = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.05
        add_highlighted_code_line(paragraph, line, 12)


def add_highlighted_code_line(paragraph, line: str, font_size: int = 12) -> None:
    indent = len(line) - len(line.lstrip(" "))
    if indent:
        run = paragraph.add_run()
        run.text = " " * indent
        style_code_run(run, COLORS["code_text"], font_size)
    rest = line[indent:]
    if rest.lstrip().startswith("#"):
        run = paragraph.add_run()
        run.text = rest
        style_code_run(run, COLORS["code_comment"], font_size)
        return
    pattern = r'("[^"\\]*(?:\\.[^"\\]*)*"|\'[^\'\\]*(?:\\.[^\'\\]*)*\'|#.*|\b[A-Za-z_][A-Za-z0-9_]*\b|\.[A-Za-z_][A-Za-z0-9_]*)'
    pos = 0
    for match in re.finditer(pattern, rest):
        if match.start() > pos:
            run = paragraph.add_run()
            run.text = rest[pos:match.start()]
            style_code_run(run, COLORS["code_text"], font_size)
        token = match.group(0)
        run = paragraph.add_run()
        run.text = token
        if token.startswith("#"):
            color = COLORS["code_comment"]
        elif token.startswith(('"', "'")):
            color = COLORS["code_string"]
        elif token in KEYWORDS:
            color = COLORS["code_keyword"]
        elif token.startswith("."):
            color = COLORS["code_func"]
        else:
            color = COLORS["code_text"]
        style_code_run(run, color, font_size, bold=(token in {"import", "from", "as"}))
        pos = match.end()
    if pos < len(rest):
        run = paragraph.add_run()
        run.text = rest[pos:]
        style_code_run(run, COLORS["code_text"], font_size)


def style_code_run(run, color: RGBColor, font_size: int, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold


def enlarge_top_right_logos(prs: Presentation, scale: float = 2.0, max_width_inches: float = 1.25) -> None:
    slide_w = prs.slide_width
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            is_top = shape.top <= Inches(0.75)
            is_right = shape.left >= slide_w - Inches(2.2)
            is_smallish = shape.width <= Inches(1.2)
            if is_top and is_right and is_smallish:
                right = shape.left + shape.width
                new_w = min(int(shape.width * scale), Inches(max_width_inches))
                ratio = shape.height / shape.width if shape.width else 1
                new_h = int(new_w * ratio)
                shape.left = right - new_w
                shape.width = new_w
                shape.height = new_h


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input", help="Input .pptx path")
    parser.add_argument("output", help="Output .pptx path")
    parser.add_argument("--logo-scale", type=float, default=1.0, help="Scale top-right logo, e.g. 2")
    parser.add_argument("--max-logo-width", type=float, default=1.25, help="Max logo width in inches")
    args = parser.parse_args()

    prs = Presentation(args.input)
    set_all_fonts(prs)
    restyle_tables(prs)
    restyle_manual_table_headers(prs)
    restyle_code_blocks(prs)
    if args.logo_scale != 1.0:
        enlarge_top_right_logos(prs, args.logo_scale, args.max_logo_width)
    prs.save(args.output)


if __name__ == "__main__":
    main()
