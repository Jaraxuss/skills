"""
Legacy Yingdao-style PPT Builder
--------------------------------
Generate simple draft PPT decks with a white + red Yingdao-like visual style.

This script is intentionally retained for compatibility only. Do not use it as
the default route for final customer-facing decks; the skill's Markdown brief
workflow with AI-led presentation design should be used instead.

Usage:
    python yingdao_ppt_builder.py deck_config_example.json output.pptx

Dependencies:
    pip install python-pptx pillow
"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

FONT = "Microsoft YaHei"

# Brand-like palette. Tune these values if the customer's template changes.
COLORS = {
    "black": RGBColor(0x1F, 0x1F, 0x1F),
    "body": RGBColor(0x33, 0x33, 0x33),
    "muted": RGBColor(0x77, 0x77, 0x77),
    "red": RGBColor(0xF2, 0x3A, 0x4A),
    "dark_red": RGBColor(0xB8, 0x3A, 0x4A),
    "table_header": RGBColor(0xF1, 0x9B, 0xAA),
    "pink": RGBColor(0xFF, 0xF3, 0xF5),
    "pink2": RGBColor(0xFF, 0xE8, 0xEC),
    "light_gray": RGBColor(0xF6, 0xF6, 0xF6),
    "line": RGBColor(0xEA, 0xEA, 0xEA),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "code_bg": RGBColor(0xFF, 0xF6, 0xF7),
    "code_border": RGBColor(0xEE, 0xA3, 0xAE),
    "code_text": RGBColor(0x40, 0x36, 0x3A),
    "code_keyword": RGBColor(0xCF, 0x2F, 0x43),
    "code_func": RGBColor(0x7E, 0x3F, 0x98),
    "code_string": RGBColor(0xB9, 0x59, 0x00),
    "code_comment": RGBColor(0x7B, 0x7B, 0x7B),
}

KEYWORDS = {
    "import", "from", "as", "for", "if", "elif", "else", "return", "in", "not",
    "and", "or", "True", "False", "None", "with", "def", "class", "lambda",
}


def cm(value: float) -> float:
    return value / 2.54


class YingdaoDeck:
    def __init__(self, logo_path: str | Path | None = None, footer: str = "From human doing to human being."):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.logo_path = Path(logo_path) if logo_path else None
        self.footer = footer
        self.page_no = 0
        self.notes_to_inject = {}

    @property
    def w(self):
        return self.prs.slide_width

    @property
    def h(self):
        return self.prs.slide_height

    def save(self, path: str | Path) -> None:
        self.prs.save(str(path))

    def set_speaker_notes(self, slide, notes_text: str | None) -> None:
        if not notes_text:
            return
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_text
        except Exception:
            self.notes_to_inject[self.page_no] = notes_text

    # ---------- base layout ----------
    def add_slide_base(self, show_logo: bool = True, show_footer: bool = True):
        slide = self.prs.slides.add_slide(self.blank)
        self.page_no += 1
        self._add_background(slide)
        if show_logo:
            self._add_logo(slide)
        if show_footer:
            self._add_footer(slide, self.page_no)
        return slide

    def _add_background(self, slide) -> None:
        # white canvas
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = COLORS["white"]

        # soft brand ambience: pale pink ring/blocks on right and corner waves
        ring = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(9.35), Inches(1.1), Inches(3.65), Inches(3.65),
        )
        ring.fill.solid()
        ring.fill.fore_color.rgb = RGBColor(0xFF, 0xEA, 0xED)
        ring.fill.transparency = 35
        ring.line.fill.background()

        cut = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(9.85), Inches(1.58), Inches(2.75), Inches(2.75),
        )
        cut.fill.solid()
        cut.fill.fore_color.rgb = COLORS["white"]
        cut.line.fill.background()

        # bottom line
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.58), Inches(7.03), Inches(12.1), Pt(1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["line"]
        line.line.fill.background()

    def _add_logo(self, slide) -> None:
        if self.logo_path and self.logo_path.exists():
            # top-right anchor; width around 0.62 inch gives visible but not dominant branding
            width = Inches(0.72)
            slide.shapes.add_picture(str(self.logo_path), Inches(11.95), Inches(0.33), width=width)
        else:
            self.text(slide, "影刀", 11.9, 0.28, 0.65, 0.25, size=14, bold=True, color=COLORS["black"])

    def _add_footer(self, slide, page_no: int) -> None:
        self.text(slide, self.footer, 0.58, 7.12, 3.2, 0.18, size=8, color=RGBColor(0x9A, 0x9A, 0x9A))
        self.text(slide, f"{page_no:02d}", 12.10, 7.10, 0.5, 0.2, size=8, color=COLORS["dark_red"], align=PP_ALIGN.RIGHT)

    def title_bar(self, slide, title: str, subtitle: str | None = None) -> None:
        mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(0.62), Inches(0.05), Inches(0.48))
        mark.fill.solid()
        mark.fill.fore_color.rgb = COLORS["red"]
        mark.line.fill.background()
        self.text(slide, title, 0.76, 0.55, 8.9, 0.48, size=24, bold=True, color=COLORS["black"])
        if subtitle:
            self.text(slide, subtitle, 0.76, 1.09, 8.9, 0.28, size=15, color=COLORS["red"], bold=True)

    # ---------- primitives ----------
    def text(self, slide, text: str, x: float, y: float, w: float, h: float, *,
             size: int = 14, bold: bool = False, color: RGBColor | None = None,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        tf.vertical_anchor = valign
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or COLORS["body"]
        return box

    def card(self, slide, x: float, y: float, w: float, h: float, *, fill=None, line=None, shadow=False):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill or COLORS["pink"]
        shp.line.color.rgb = line or RGBColor(0xF1, 0xC9, 0xD0)
        shp.line.width = Pt(0.75)
        if shadow:
            # python-pptx has limited shadow support; use light border instead for compatibility.
            pass
        return shp

    def labeled_card(self, slide, title: str, body: str, x: float, y: float, w: float, h: float,
                     *, accent=True, body_size=14):
        self.card(slide, x, y, w, h)
        if accent:
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y + 0.18), Inches(0.055), Inches(h - 0.36))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLORS["red"]
            bar.line.fill.background()
        self.text(slide, title, x + 0.24, y + 0.18, w - 0.38, 0.32, size=16, bold=True, color=COLORS["black"])
        self.text(slide, body, x + 0.24, y + 0.62, w - 0.38, h - 0.78, size=body_size, color=COLORS["body"], valign=MSO_ANCHOR.TOP)

    def arrow(self, slide, x: float, y: float, w: float, h: float):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS["red"]
        shp.line.fill.background()
        return shp

    # ---------- content layouts ----------
    def add_cover(self, cfg: dict) -> None:
        slide = self.add_slide_base(show_logo=True, show_footer=True)
        # optional hero image on right
        hero = cfg.get("hero_image")
        if hero and Path(hero).exists():
            slide.shapes.add_picture(hero, Inches(6.6), Inches(1.25), width=Inches(6.2))
        self.text(slide, cfg.get("eyebrow", "客户培训"), 0.75, 1.0, 5.2, 0.35, size=15, color=COLORS["red"], bold=True)
        self.text(slide, cfg["title"], 0.75, 1.55, 5.6, 1.6, size=34, color=COLORS["black"], bold=True, valign=MSO_ANCHOR.TOP)
        if cfg.get("subtitle"):
            self.text(slide, cfg["subtitle"], 0.78, 3.15, 5.4, 0.8, size=16, color=COLORS["body"], valign=MSO_ANCHOR.TOP)
        if cfg.get("meta"):
            self.card(slide, 0.78, 4.55, 5.1, 0.78, fill=COLORS["pink2"])
            self.text(slide, cfg["meta"], 1.0, 4.72, 4.65, 0.45, size=14, color=COLORS["body"])
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def add_big_statement(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.title_bar(slide, cfg.get("section", "核心观点"), cfg.get("subtitle"))
        self.text(slide, cfg["statement"], 1.2, 2.35, 10.8, 1.3, size=30, bold=True, color=COLORS["black"], align=PP_ALIGN.CENTER)
        if cfg.get("note"):
            self.text(slide, cfg["note"], 1.65, 4.0, 10.0, 0.55, size=16, color=COLORS["red"], bold=True, align=PP_ALIGN.CENTER)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def add_table_slide(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.title_bar(slide, cfg["title"], cfg.get("subtitle"))
        rows = cfg["rows"]
        headers = cfg.get("headers")
        table_rows = len(rows) + (1 if headers else 0)
        table_cols = len(headers or rows[0])
        x, y, w, h = cfg.get("bbox", [0.76, 1.55, 11.7, 4.9])
        table_shape = slide.shapes.add_table(table_rows, table_cols, Inches(x), Inches(y), Inches(w), Inches(h))
        table = table_shape.table
        widths = cfg.get("col_widths")
        if widths:
            total = sum(widths)
            for i, rel in enumerate(widths):
                table.columns[i].width = Inches(w * rel / total)
        row_idx = 0
        if headers:
            for c, header in enumerate(headers):
                self._format_cell(table.cell(0, c), header, fill=COLORS["table_header"], bold=True, size=14, color=COLORS["black"])
            row_idx = 1
        for r, row in enumerate(rows, start=row_idx):
            for c, value in enumerate(row):
                fill = COLORS["white"] if (r - row_idx) % 2 == 0 else COLORS["pink"]
                self._format_cell(table.cell(r, c), str(value), fill=fill, bold=False, size=14, color=COLORS["body"])
        if cfg.get("footer_note"):
            self.text(slide, cfg["footer_note"], 0.95, 6.36, 10.8, 0.35, size=14, color=COLORS["body"], align=PP_ALIGN.CENTER)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def _format_cell(self, cell, text: str, *, fill, bold: bool, size: int, color):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        tf.clear()
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if bold else PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def add_cards_slide(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.title_bar(slide, cfg["title"], cfg.get("subtitle"))
        cards = cfg["cards"]
        n = len(cards)
        gap = 0.28
        total_w = 11.3
        card_w = (total_w - gap * (n - 1)) / n
        x0 = 1.0
        y = cfg.get("y", 1.65)
        h = cfg.get("card_height", 3.6)
        for i, c in enumerate(cards):
            x = x0 + i * (card_w + gap)
            self.labeled_card(slide, c["title"], c.get("body", ""), x, y, card_w, h, body_size=c.get("body_size", 14))
            if c.get("tag"):
                self.text(slide, c["tag"], x + card_w - 1.15, y + 0.2, 0.85, 0.22, size=11, color=COLORS["red"], align=PP_ALIGN.RIGHT)
        if cfg.get("bottom_note"):
            self.text(slide, cfg["bottom_note"], 1.15, 5.75, 10.8, 0.6, size=16, bold=True, color=COLORS["black"], align=PP_ALIGN.CENTER)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def add_process_slide(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.title_bar(slide, cfg["title"], cfg.get("subtitle"))
        steps = cfg["steps"]
        x0, y, w, h, gap = 0.85, 2.0, 1.65, 1.12, 0.37
        for i, step in enumerate(steps):
            x = x0 + i * (w + gap)
            self.card(slide, x, y, w, h, fill=COLORS["pink"])
            self.text(slide, step, x + 0.12, y + 0.22, w - 0.24, 0.6, size=14, color=COLORS["black"], bold=True, align=PP_ALIGN.CENTER)
            if i < len(steps) - 1:
                self.arrow(slide, x + w + 0.08, y + 0.37, 0.24, 0.24)
        if cfg.get("notes"):
            x_cards = [0.95, 4.75, 8.55]
            for i, note in enumerate(cfg["notes"][:3]):
                self.labeled_card(slide, note["title"], note["body"], x_cards[i], 4.2, 3.35, 1.25, body_size=13)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def add_code_case_slide(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.title_bar(slide, cfg["title"], cfg.get("subtitle"))
        # left business cards
        for i, card in enumerate(cfg.get("cards", [])):
            self.labeled_card(slide, card["title"], card.get("body", ""), 0.82, 1.55 + i * 1.42, 3.6, 1.18, body_size=13)
        # code block
        self.add_code_block(slide, cfg["code"], 4.75, 1.48, 7.25, 4.3, font_size=12)
        if cfg.get("flow"):
            steps = cfg["flow"]
            step_w = 1.42
            x0 = 1.0
            y = 6.0
            for i, step in enumerate(steps):
                self.text(slide, step, x0 + i * 2.05, y, step_w, 0.28, size=13, color=COLORS["black"], bold=True, align=PP_ALIGN.CENTER)
                if i < len(steps) - 1:
                    self.arrow(slide, x0 + i * 2.05 + 1.52, y + 0.03, 0.25, 0.2)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))

    def add_code_block(self, slide, code: str, x: float, y: float, w: float, h: float, font_size: int = 12):
        bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        bg.name = "CODE_BACKGROUND"
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLORS["code_bg"]
        bg.line.color.rgb = COLORS["code_border"]
        bg.line.width = Pt(0.8)
        box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.16), Inches(w - 0.32), Inches(h - 0.32))
        box.name = "CODE_TEXT"
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        for i, line in enumerate(code.strip("\n").split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(0)
            p.line_spacing = 1.05
            self._add_highlighted_code_line(p, line, font_size)
        return box

    def _add_highlighted_code_line(self, paragraph, line: str, font_size: int = 12):
        # Preserve leading spaces by adding them as a run.
        indent = len(line) - len(line.lstrip(" "))
        if indent:
            run = paragraph.add_run()
            run.text = " " * indent
            self._style_run(run, COLORS["code_text"], font_size)
        rest = line[indent:]
        if rest.lstrip().startswith("#"):
            run = paragraph.add_run()
            run.text = rest
            self._style_run(run, COLORS["code_comment"], font_size)
            return
        # Tokenize strings, comments and common identifiers.
        pattern = r'("[^"\\]*(?:\\.[^"\\]*)*"|\'[^\'\\]*(?:\\.[^\'\\]*)*\'|#.*|\b[A-Za-z_][A-Za-z0-9_]*\b|\.[A-Za-z_][A-Za-z0-9_]*)'
        pos = 0
        for m in re.finditer(pattern, rest):
            if m.start() > pos:
                run = paragraph.add_run()
                run.text = rest[pos:m.start()]
                self._style_run(run, COLORS["code_text"], font_size)
            token = m.group(0)
            run = paragraph.add_run()
            run.text = token
            if token.startswith("#"):
                color = COLORS["code_comment"]
            elif token.startswith(('"', "'")):
                color = COLORS["code_string"]
            elif token in KEYWORDS:
                color = COLORS["code_keyword"]
            elif token.startswith("."):
                color = COLORS["code_func"]
            else:
                color = COLORS["code_text"]
            self._style_run(run, color, font_size, bold=(token in {"import", "from", "as"}))
            pos = m.end()
        if pos < len(rest):
            run = paragraph.add_run()
            run.text = rest[pos:]
            self._style_run(run, COLORS["code_text"], font_size)

    def _style_run(self, run, color: RGBColor, font_size: int = 12, bold: bool = False):
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold

    def add_qa(self, cfg: dict) -> None:
        slide = self.add_slide_base()
        self.text(slide, cfg.get("statement", "Q&A"), 1.0, 2.05, 11.2, 0.9, size=30, bold=True, color=COLORS["black"], align=PP_ALIGN.CENTER)
        if cfg.get("note"):
            self.text(slide, cfg["note"], 1.4, 3.2, 10.5, 0.55, size=17, bold=True, color=COLORS["red"], align=PP_ALIGN.CENTER)
        self.set_speaker_notes(slide, cfg.get("speaker_notes"))


def build_from_config(config_path: str | Path, output_path: str | Path) -> None:
    cfg = json.loads(Path(config_path).read_text(encoding="utf-8"))
    base = Path(config_path).parent
    logo_path = cfg.get("logo_path", "assets/yingdao_logo.png")
    logo_path = str((base / logo_path).resolve()) if logo_path else None
    deck = YingdaoDeck(logo_path=logo_path, footer=cfg.get("footer", "From human doing to human being."))
    for slide_cfg in cfg["slides"]:
        slide_type = slide_cfg["type"]
        if slide_type == "cover":
            if slide_cfg.get("hero_image"):
                slide_cfg = dict(slide_cfg)
                slide_cfg["hero_image"] = str((base / slide_cfg["hero_image"]).resolve())
            deck.add_cover(slide_cfg)
        elif slide_type == "statement":
            deck.add_big_statement(slide_cfg)
        elif slide_type == "table":
            deck.add_table_slide(slide_cfg)
        elif slide_type == "cards":
            deck.add_cards_slide(slide_cfg)
        elif slide_type == "process":
            deck.add_process_slide(slide_cfg)
        elif slide_type == "code_case":
            deck.add_code_case_slide(slide_cfg)
        elif slide_type == "qa":
            deck.add_qa(slide_cfg)
        else:
            raise ValueError(f"Unsupported slide type: {slide_type}")
    deck.save(output_path)
    if getattr(deck, "notes_to_inject", None):
        inject_speaker_notes_to_pptx(Path(output_path), deck.notes_to_inject)


# ---------- Fallback Speaker Notes XML Injection Utilities ----------
REL_NOTES_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
REL_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
CT_NOTES_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"


def _next_relationship_id(rels_xml: str) -> str:
    ids = [int(match) for match in re.findall(r'Id="rId(\d+)"', rels_xml)]
    return f"rId{max(ids, default=0) + 1}"


def _ensure_content_type(content_types_xml: str, slide_number: int) -> str:
    part_name = f"/ppt/notesSlides/notesSlide{slide_number}.xml"
    if part_name in content_types_xml:
        return content_types_xml

    override = (
        f'<Override PartName="{part_name}" '
        f'ContentType="{CT_NOTES_SLIDE}"/>'
    )
    return content_types_xml.replace("</Types>", f"{override}</Types>")


def _ensure_slide_to_notes_relationship(rels_xml: str, slide_number: int) -> str:
    target = f"../notesSlides/notesSlide{slide_number}.xml"
    if REL_NOTES_SLIDE in rels_xml and target in rels_xml:
        return rels_xml

    rel_id = _next_relationship_id(rels_xml)
    relationship = (
        f'<Relationship Id="{rel_id}" '
        f'Type="{REL_NOTES_SLIDE}" '
        f'Target="{target}"/>'
    )
    return rels_xml.replace("</Relationships>", f"{relationship}</Relationships>")


def _notes_paragraphs_xml(notes_text: str) -> str:
    from xml.sax.saxutils import escape
    paragraphs = []
    for line in notes_text.splitlines() or [""]:
        safe_line = escape(line)
        paragraphs.append(
            "<a:p>"
            "<a:r>"
            "<a:rPr lang=\"zh-CN\" sz=\"1500\" dirty=\"0\"/>"
            f"<a:t>{safe_line}</a:t>"
            "</a:r>"
            "<a:endParaRPr lang=\"zh-CN\" sz=\"1500\"/>"
            "</a:p>"
        )
    return "".join(paragraphs)


def _notes_slide_xml(notes_text: str) -> str:
    paragraphs = _notes_paragraphs_xml(notes_text)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
         xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
         xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/><a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>
        <p:nvSpPr>
          <p:cNvPr id="2" name="Notes Placeholder 1"/>
          <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>
          <p:nvPr><p:ph type="body" idx="1"/></p:nvPr>
        </p:nvSpPr>
        <p:spPr>
          <a:xfrm><a:off x="685800" y="457200"/><a:ext cx="5486400" cy="8229600"/></a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        </p:spPr>
        <p:txBody>
          <a:bodyPr wrap="square" rtlCol="0"><a:spAutoFit/></a:bodyPr>
          <a:lstStyle/>
          {paragraphs}
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''


def _notes_slide_rels_xml(slide_number: int) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="{REL_SLIDE}" Target="../slides/slide{slide_number}.xml"/>
</Relationships>'''


def inject_speaker_notes_to_pptx(pptx_path: Path, slide_notes: dict[int, str]) -> None:
    """Inject speaker notes for multiple slides into a .pptx file via ZIP manipulation."""
    import zipfile
    import shutil
    import tempfile

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir) / "work.pptx"
        shutil.copy2(pptx_path, tmp_path)

        with zipfile.ZipFile(tmp_path, "r") as zin:
            parts = {name: zin.read(name) for name in zin.namelist()}

        for slide_number, notes_text in slide_notes.items():
            slide_xml = f"ppt/slides/slide{slide_number}.xml"
            slide_rels_xml = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
            notes_slide_xml = f"ppt/notesSlides/notesSlide{slide_number}.xml"
            notes_slide_rels_xml = f"ppt/notesSlides/_rels/notesSlide{slide_number}.xml.rels"

            if slide_xml not in parts:
                continue

            parts["[Content_Types].xml"] = _ensure_content_type(
                parts["[Content_Types].xml"].decode("utf-8"), slide_number
            ).encode("utf-8")

            if slide_rels_xml not in parts:
                parts[slide_rels_xml] = b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>'

            parts[slide_rels_xml] = _ensure_slide_to_notes_relationship(
                parts[slide_rels_xml].decode("utf-8"), slide_number
            ).encode("utf-8")
            parts[notes_slide_xml] = _notes_slide_xml(notes_text).encode("utf-8")
            parts[notes_slide_rels_xml] = _notes_slide_rels_xml(slide_number).encode("utf-8")

        with zipfile.ZipFile(pptx_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for name, data in parts.items():
                zout.writestr(name, data)


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python yingdao_ppt_builder.py deck_config.json output.pptx")
        sys.exit(1)
    build_from_config(sys.argv[1], sys.argv[2])
